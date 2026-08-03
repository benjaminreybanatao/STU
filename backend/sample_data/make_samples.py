"""Generates synthetic OM PDF/PPTX + underwriting model fixtures for
end-to-end testing. Not part of the app itself -- just test fixtures."""
import fitz
from PIL import Image, ImageDraw
import openpyxl
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

OUT = Path(__file__).parent

# --- Fake "hero photo" (large, photographic aspect ratio) ---
hero = Image.new("RGB", (1600, 1000), (90, 110, 140))
d = ImageDraw.Draw(hero)
for i in range(0, 1600, 40):
    d.line([(i, 0), (i, 1000)], fill=(80 + i % 60, 100, 130), width=20)
d.rectangle([100, 600, 1500, 950], fill=(60, 60, 65))
hero.save(OUT / "hero.jpg", quality=85)

# --- Fake logo (thin banner, should be filtered out) ---
logo = Image.new("RGB", (400, 60), (255, 255, 255))
d2 = ImageDraw.Draw(logo)
d2.text((10, 20), "ACME BROKERAGE LOGO", fill=(0, 0, 0))
logo.save(OUT / "logo.png")

# --- Fake floor plan (small, near-square line art) ---
plan = Image.new("RGB", (250, 250), (245, 245, 245))
d3 = ImageDraw.Draw(plan)
d3.rectangle([20, 20, 230, 230], outline=(0, 0, 0), width=2)
plan.save(OUT / "floorplan.png")

doc = fitz.open()
page = doc.new_page(width=612, height=792)

text = (
    "1200 Market Street, Austin, TX 78701\n\n"
    "Offering Memorandum\n\n"
    "1200 Market Street is a 145,000 square feet office property built in 2016, "
    "located in the Austin CBD submarket, a high-growth technology submarket with "
    "strong absorption over the trailing 24 months.\n\n"
    "The property is 94% leased. Tenants include a mix of technology and professional "
    "services firms anchored by a 12-year investment-grade tenant occupying 40% of the NRA.\n\n"
)
page.insert_textbox(fitz.Rect(50, 50, 560, 300), text, fontsize=11)

page.insert_image(fitz.Rect(50, 320, 560, 620), filename=str(OUT / "hero.jpg"))
page.insert_image(fitz.Rect(50, 630, 200, 660), filename=str(OUT / "logo.png"))

page2 = doc.new_page(width=612, height=792)
page2.insert_textbox(fitz.Rect(50, 50, 560, 100), "Site Plan", fontsize=14)
page2.insert_image(fitz.Rect(180, 120, 430, 370), filename=str(OUT / "floorplan.png"))

doc.save(OUT / "sample_om.pdf")
doc.close()

# --- Same OM content, as a PPTX (brokers send both formats) ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(3))
box.text_frame.text = (
    "900 Congress Avenue, Austin, TX 78701 Offering Memorandum. "
    "900 Congress Avenue is a 210,000 square feet office building built in 2019, "
    "located in the Austin CBD submarket, a dense urban submarket. "
    "The property is 88% leased. Tenants include several technology tenants."
)
slide.shapes.add_picture(str(OUT / "hero.jpg"), Inches(0.5), Inches(4), Inches(6), Inches(3))
slide.shapes.add_picture(str(OUT / "logo.png"), Inches(7), Inches(4), Inches(2), Inches(0.3))
prs.save(OUT / "sample_om.pptx")

# --- Underwriting model ---
# Laid out the way the real firm models are: a label column, a "Gross" column
# and an explicit "$ Per Unit" column, plus a Sources & Uses block off to the
# right. The parser reads per-unit figures only from that stated column -- it
# never derives them -- so the fixture needs them present to exercise it.
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Base Case"

# The OM fixture is a 145,000 SF office deal, so the model states figures
# per square foot. "$ Per Unit" is just the column's heading in these models.
SF = 145_000

rows = [
    # (label, gross, per-unit)
    ("Building Name", "1200 Market Street", None),
    ("Total Square Feet", SF, None),
    ("Current Occupancy", 0.94, None),
    ("Occupancy at Exit", 1.00, None),
    ("Projected Hold Period", "5 Years", None),
    ("Pricing", "Gross", "$ Per Unit"),
    ("Purchase Price", 42_500_000, 293),
    ("Peak Cost", 44_100_000, 304),
    ("Exit Price (Gross)", 57_800_000, 399),
    ("Cap Rate", None, None),
    ("Going-In Cap Rate", None, 0.0525),
    ("Market Cap Rate", None, 0.0560),
    ("Exit Cap Rate", None, 0.0575),
    ("Gross Returns", None, None),
    ("Unlevered IRR", None, 0.079),
    ("Unlevered Equity Multiple", None, 1.42),
    ("Levered IRR", None, 0.134),
    ("Levered Equity Multiple", None, 1.71),
    ("Cash-on-Cash", None, 0.061),
    ("Debt", "Gross", None),
    ("Leverage", 0.60, 25_500_000),
    ("Interest Rate", 0.055, None),
    ("Equity", "Gross", "$ Per Unit"),
    ("Initial Equity", 17_000_000, 117),
    ("Peak Equity", 17_850_000, 123),
    ("Lease Term Assumption", "10-year weighted average", None),
    ("Downtime Assumption", "9 months between leases", None),
    ("Exit Assumption", "Sale in Year 5 at stabilized NOI", None),
]
for i, (label, gross, per_unit) in enumerate(rows, start=2):
    ws.cell(row=i, column=1, value=label)
    if gross is not None:
        ws.cell(row=i, column=2, value=gross)
    if per_unit is not None:
        ws.cell(row=i, column=3, value=per_unit)

# Sources & Uses at Close, in its own block to the right.
su_rows = [
    ("Total Sources", "Total", "$ Per Unit"),
    ("Equity", 17_000_000, 117),
    ("Gross Debt Proceeds", 25_500_000, 176),
    ("Total Sources", 42_500_000, 293),
    (None, None, None),
    ("Total Uses", "Total", "$ Per Unit"),
    ("Purchase Price", 42_500_000, 293),
    ("DD / Closing Costs", 420_000, 3),
    ("Working Capital", 300_000, 2),
    ("Equity Subtotal", 43_220_000, 298),
    ("Financing Cost", 255_000, 2),
    ("Total Uses", 43_475_000, 300),
]
for i, (label, total, per_unit) in enumerate(su_rows, start=2):
    if label is None:
        continue
    ws.cell(row=i, column=6, value=label)
    if total is not None:
        ws.cell(row=i, column=7, value=total)
    if per_unit is not None:
        ws.cell(row=i, column=8, value=per_unit)

ws2 = wb.create_sheet("Downside Case")
downside_rows = [
    ("Exit Cap Rate", None, 0.065),
    ("Levered IRR", None, 0.081),
    ("Unlevered IRR", None, 0.052),
    ("Levered Equity Multiple", None, 1.32),
    ("Unlevered Equity Multiple", None, 1.21),
    ("Cash-on-Cash", None, 0.041),
    ("Exit Price (Gross)", 49_300_000, 340),
]
for i, (label, gross, per_unit) in enumerate(downside_rows, start=2):
    ws2.cell(row=i, column=1, value=label)
    if gross is not None:
        ws2.cell(row=i, column=2, value=gross)
    if per_unit is not None:
        ws2.cell(row=i, column=3, value=per_unit)

wb.save(OUT / "sample_model.xlsx")

print("Wrote sample_om.pdf, sample_om.pptx, and sample_model.xlsx to", OUT)
