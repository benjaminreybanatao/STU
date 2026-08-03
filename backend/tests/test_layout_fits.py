"""Layout guard for the generated deck.

The exhibit tables are laid out at measured absolute coordinates with cell
wrapping switched off, which keeps them faithful to the reference decks but
means over-long text silently spills outside its column instead of reflowing.
That is exactly how the Building Name cell once pushed a full postal address
across the gap into the Sources & Uses table.

These tests assert every shape stays on the slide and every table cell's text
fits the width it was given, for both a full deal and an OM-only deal (where
every financial cell is a placeholder).

Run with: python -m pytest tests/ -q     (from the backend/ directory)
"""
import itertools
import tempfile
from pathlib import Path

import pytest
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from app.pipeline import build_deal, generate_deck
from app.schema import Deal

EMU_PER_PT = 12700
SLIDE_W_PT = 960.0
SLIDE_H_PT = 540.0
CELL_INSET_PT = 10.0  # 5pt left + 5pt right, per _style_cell
FOOTER_TOP_PT = 514.0  # measured footer baseline; tables must stay above it

SAMPLE_DIR = Path(__file__).resolve().parent.parent / "sample_data"

# DejaVu is metrically wider than Gandhi Sans / Calibri, so measuring with it
# makes this check conservative: anything that passes here has headroom in the
# real fonts.
_FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
_SCALE = 10  # measure at 10x for sub-point resolution
_measure = ImageDraw.Draw(Image.new("RGB", (8, 8)))


def _text_width_pt(text: str, size_pt: float) -> float:
    font = ImageFont.truetype(_FONT_PATH, int(round(size_pt * _SCALE)))
    return _measure.textlength(text, font=font) / _SCALE


def _build(tmp_path, with_model: bool):
    out = tmp_path / "deck.pptx"
    deal = build_deal(
        str(SAMPLE_DIR / "sample_om.pdf"),
        str(SAMPLE_DIR / "sample_model.xlsx") if with_model else None,
        {},
        str(tmp_path),
    )
    generate_deck(deal, str(out))
    return Presentation(str(out))


@pytest.fixture(scope="module")
def decks(tmp_path_factory):
    return {
        "full": _build(tmp_path_factory.mktemp("full"), True),
        "om_only": _build(tmp_path_factory.mktemp("om"), False),
    }


@pytest.mark.parametrize("variant", ["full", "om_only"])
def test_all_shapes_stay_on_slide(decks, variant):
    prs = decks[variant]
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            right = (shape.left + shape.width) / EMU_PER_PT
            bottom = (shape.top + shape.height) / EMU_PER_PT
            assert right <= SLIDE_W_PT + 0.5, f"slide {i}: {shape.shape_type} runs off the right edge"
            assert bottom <= SLIDE_H_PT + 0.5, f"slide {i}: {shape.shape_type} runs off the bottom"


def _cell_overflows(prs) -> list[str]:
    """Each paragraph (line) in a cell is checked independently, not the
    cell's full joined text -- the 3-column Transaction Overview layout
    stacks a row's Gross and $ Per Unit figures as two separate lines in one
    cell, each independently guaranteed to fit rather than a combined one
    that couldn't."""
    overflows = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            col_widths = [c.width / EMU_PER_PT for c in shape.table.columns]
            for r, row in enumerate(shape.table.rows):
                for c, cell in enumerate(row.cells):
                    if getattr(cell, "is_spanned", False):
                        continue
                    span = getattr(cell, "span_width", 1)
                    available = sum(col_widths[c:c + span]) - CELL_INSET_PT
                    for p in cell.text_frame.paragraphs:
                        line = "".join(run.text for run in p.runs).strip()
                        if not line:
                            continue
                        size_pt = next(
                            (run.font.size.pt for run in p.runs if run.font.size), 11.2
                        )
                        needed = _text_width_pt(line, size_pt)
                        if needed > available:
                            overflows.append(
                                f"slide {i} r{r}c{c}: needs {needed:.0f}pt, has {available:.0f}pt -- {line!r}"
                            )
    return overflows


@pytest.mark.parametrize("variant", ["full", "om_only"])
def test_table_cell_text_fits_its_column(decks, variant):
    overflows = _cell_overflows(decks[variant])
    assert not overflows, "table text overflows its column:\n  " + "\n  ".join(overflows)


def test_building_name_with_no_comma_does_not_wrap():
    """A bare "123 Long Street Name Boulevard" address (no city/state on the
    same field, so no comma to split on) used to pass through _building_name
    unchanged, wrap to two lines in the narrow value column, and grow that
    row -- pushing the whole table past the footer. Regression coverage for
    a real deal that hit exactly this."""
    deal = Deal()
    deal.om_facts["address"] = "2390 Mission College Boulevard"
    deal.model_facts.update({
        "sf_or_units": 153549, "occupancy": 0.46, "occupancy_at_exit": 1.0, "hold_period": 4,
        "purchase_price": 41000000, "peak_cost": 67211746, "exit_price": 46991547,
        "going_in_cap": 0.041, "market_cap": 0.08, "exit_cap": 0.08,
        "unlevered_irr": 0.146, "unlevered_em": 1.65, "levered_irr": 0.208, "levered_em": 1.96,
        "leverage": 1.718, "gross_debt_proceeds": 26650000,
        "initial_equity": 16245688, "peak_equity": 24393502,
    })
    deal._downside_model_facts = {}
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "deck.pptx"
        generate_deck(deal, str(out))
        prs = Presentation(str(out))

    assert not _cell_overflows(prs)


@pytest.mark.parametrize("variant", ["full", "om_only"])
def test_tables_clear_the_footer_even_if_powerpoint_grows_rows(decks, variant):
    """A declared row height is only a minimum.

    PowerPoint grows each row to fit its line and ignores wrap="none" on table
    cells, so a table sized purely on declared heights can still run off the
    slide once the brand font is substituted. Re-check the bottom edge against
    the height PowerPoint would actually use.
    """
    # Deliberately NOT imported from the builder: this is an external fact
    # about PowerPoint. 1.58 was measured from real output (11.2pt text on a
    # ~17.7pt row with the brand font substituted) and still proved too low, so
    # this asserts the stronger margin the builder now designs for. Importing
    # the builder's own constant would make the test merely self-consistent and
    # would stop catching a regression that changed both together.
    POWERPOINT_MIN_ROW_PER_PT = 1.78

    footer_top_pt = FOOTER_TOP_PT
    for i, slide in enumerate(decks[variant].slides, start=1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            grown = 0.0
            for r, row in enumerate(table.rows):
                declared = row.height / EMU_PER_PT
                # A cell with 2 stacked lines (the 3-column layout's Gross /
                # $ Per Unit pair) needs roughly twice a single line's grown
                # height, not just the bigger font -- summed within a cell
                # (its lines stack), then the worst cell wins for the row
                # (cells in the same row sit side by side, not stacked).
                needed = 0.0
                for c in range(len(table.columns)):
                    cell = table.cell(r, c)
                    if getattr(cell, "is_spanned", False):
                        continue
                    cell_needed = sum(
                        max((run.font.size.pt for run in p.runs if run.font.size), default=0.0)
                        * POWERPOINT_MIN_ROW_PER_PT
                        for p in cell.text_frame.paragraphs
                        if any(run.text.strip() for run in p.runs)
                    )
                    needed = max(needed, cell_needed)
                grown += max(declared, needed)
            bottom = shape.top / EMU_PER_PT + grown
            assert bottom <= footer_top_pt, (
                f"slide {i}: table would render {bottom:.1f}pt deep, "
                f"past the footer at {footer_top_pt:.1f}pt"
            )


@pytest.mark.parametrize("variant", ["full", "om_only"])
def test_exhibit_tables_do_not_collide(decks, variant):
    """No two tables on the exhibit slide overlap.

    The Transaction Overview header/mini tables and Sources & Uses only need
    checking against tables they actually share vertical space with -- the
    header spans the same x-range as its 3 mini tables below it, which is
    fine since they're stacked, not overlapping.
    """
    prs = decks[variant]
    for i, slide in enumerate(prs.slides, start=1):
        boxes = [
            (sh.left / EMU_PER_PT, sh.top / EMU_PER_PT,
             (sh.left + sh.width) / EMU_PER_PT, (sh.top + sh.height) / EMU_PER_PT)
            for sh in slide.shapes if sh.has_table
        ]
        for (l1, t1, r1, b1), (l2, t2, r2, b2) in itertools.combinations(boxes, 2):
            vertically_overlaps = t1 < b2 and t2 < b1
            horizontally_overlaps = l1 < r2 and l2 < r1
            assert not (vertically_overlaps and horizontally_overlaps), (
                f"slide {i}: two tables overlap "
                f"(({l1:.1f},{t1:.1f})-({r1:.1f},{b1:.1f}) vs "
                f"({l2:.1f},{t2:.1f})-({r2:.1f},{b2:.1f}))"
            )
