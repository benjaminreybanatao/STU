"""OM parser regression coverage (text-fact extraction and image extraction).

Run with: python -m pytest tests/ -q     (from the backend/ directory)
"""
import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.parsing.om_parser import _extract_images_from_pptx, _facts_from_text, _normalize_to_rgb

FIXTURES_DIR = Path(__file__).resolve().parent / "fixtures"


def test_occupancy_does_not_cross_a_line_break():
    """A real OM PDF's text layer had "...163,545\\n% Leased\\n100%\\nWALT" --
    a size figure ending in "545", a "% Leased" label on the next line, and
    (elsewhere, further down) the real value in "100% occupancy". \\s* in the
    old pattern matches newlines, so it matched "545" + the line break +
    "% Leased" -- reporting 545% occupancy -- before .search() ever reached
    the real, same-line match. Regression coverage using the exact text shape
    that triggered it, not the real PDF (this is pure regex behavior)."""
    text = (
        "As-Leased | BOMA\n139,661 | 163,545\n% Leased\n100%\nWALT Rem. (Aug\n"
        "flow, supported by 100% occupancy and\nover $10 million"
    )
    assert _facts_from_text(text).get("occupancy") == "100%"


def test_occupancy_same_line_still_matches():
    """The common, same-line phrasing this pattern targets must keep working."""
    assert _facts_from_text("The Property is 97% leased as of today.").get("occupancy") == "97%"

# A 64x64 downsize of a real aerial-map photo from an actual OM PDF (a
# DeviceCMYK JPEG, Adobe APP14 marker, transform byte 0 -- the common case for
# a PDF produced by an Acrobat Distiller/InDesign print pipeline). Kept as a
# real fixture rather than a synthetic Pillow-encoded CMYK JPEG on purpose:
# Pillow's own CMYK JPEG encoder round-trips fine through Pillow's own
# decoder, which makes a synthetic test blind to this exact bug. Real
# Adobe-pipeline exports don't use Pillow's convention, so only real fixture
# bytes actually exercise it.
REAL_CMYK_SAMPLE = FIXTURES_DIR / "real_cmyk_sample.jpg"


def test_real_cmyk_jpeg_is_not_a_color_negative():
    """Pillow's JPEG decoder unconditionally treats every CMYK-mode JPEG as
    "Adobe inverted" and un-inverts it while decoding -- but real photography
    embedded in OM PDFs isn't actually stored that way, so a plain
    `.convert("RGB")` renders it as a color negative (near-black here; this
    is a green, sunlit aerial map). A previous version of this function
    tried to detect the few cases needing correction via the JPEG APP14
    marker's transform byte, which was backwards: this real file uses
    transform 0, not the transform 2 that version treated as the inverted
    case, so it was never actually corrected."""
    normalized = _normalize_to_rgb(REAL_CMYK_SAMPLE.read_bytes())
    assert normalized.mode == "RGB"
    r, g, b = normalized.getpixel((32, 32))
    # The broken (non-inverted) decode of this exact pixel is near-black,
    # (24, 20, 25) -- nowhere near this range.
    assert (r, g, b) > (60, 60, 60), f"looks like a color negative: {(r, g, b)}"


def test_unloadable_picture_is_skipped_not_fatal(tmp_path):
    """A picture that Image.open() accepts (valid header) but can't actually
    be decoded (truncated data) used to crash extraction for the whole deck,
    since the failure only surfaces on save()/load(), after the original
    try/except around open() had already passed. One bad picture should be
    skipped, not take down every other picture in the deck."""
    truncated_path = tmp_path / "truncated.png"
    good_path = tmp_path / "good.png"

    full = io.BytesIO()
    Image.new("RGB", (400, 300), (10, 20, 30)).save(full, format="PNG")
    truncated_path.write_bytes(full.getvalue()[: len(full.getvalue()) // 2])
    Image.new("RGB", (400, 300), (50, 60, 70)).save(good_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(truncated_path), Inches(0), Inches(0), width=Inches(2))
    slide.shapes.add_picture(str(good_path), Inches(3), Inches(0), width=Inches(2))
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    out_dir = tmp_path / "extracted"
    out_dir.mkdir()
    results = _extract_images_from_pptx(str(pptx_path), str(out_dir))

    assert len(results) == 1
    assert Image.open(results[0]["path"]).getpixel((0, 0)) == (50, 60, 70)
