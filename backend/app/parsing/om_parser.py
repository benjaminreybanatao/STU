"""
OM (offering memorandum) parser -- accepts either a PDF or a PowerPoint deck,
since brokers send both.

Two jobs, kept deliberately separate:
  1. extract_facts() - regex/heuristic pull of narrative deal facts out of
     the raw text layer (address, year built, SF/units, occupancy, etc).
  2. extract_images() - pull embedded raster images out of the document and
     rank them so the "hero photo" picked for slide 1 is an actual building
     photograph, not a logo, chart, or floor-plan line-art graphic.

Both are heuristic by nature -- an OM is free-form marketing prose, not a
structured document -- so every extracted fact is a best-effort regex match.
Nothing here is invented; if a pattern doesn't match, the field is simply
left out of the returned dict and the gap analysis will flag it as missing.

PDF text/images are read via PyMuPDF; PPTX text/images are read via
python-pptx (PyMuPDF doesn't parse OOXML slide decks). Both formats funnel
into the same regex fact-extraction and the same image size/aspect filter
and RGB-normalization, so behavior is identical regardless of source format.
"""
import io
import re
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_EXTENSIONS = (".pptx", ".pptm")

ADDRESS_RE = re.compile(
    r"\b(\d{2,6}\s+[A-Z][A-Za-z0-9.'\-]*(?:\s+[A-Z][A-Za-z0-9.'\-]*){0,4}"
    r"\s+(?:Street|St|Avenue|Ave|Boulevard|Blvd|Road|Rd|Drive|Dr|Way|Lane|Ln|"
    r"Parkway|Pkwy|Court|Ct|Place|Pl|Highway|Hwy)\.?,?\s*"
    r"(?:[A-Za-z .]+,\s*)?[A-Z]{2}\s*\d{5}(?:-\d{4})?)\b"
)

YEAR_BUILT_RE = re.compile(
    r"(?:built|constructed|delivered|year built)\D{0,15}((?:19|20)\d{2})",
    re.IGNORECASE,
)

SF_RE = re.compile(
    r"([\d,]{4,10})\s*(?:rentable\s+)?(?:square feet|sf|rsf)\b", re.IGNORECASE
)
UNITS_RE = re.compile(r"([\d,]{1,5})\s*(?:units|unit count|residential units)\b", re.IGNORECASE)

OCCUPANCY_RE = re.compile(r"(\d{1,3}(?:\.\d+)?)\s*%\s*(?:leased|occupied|occupancy)", re.IGNORECASE)

PROPERTY_TYPE_KEYWORDS = {
    "multifamily": ["multifamily", "apartment", "residential units"],
    "office": ["office building", "office property", "office asset"],
    "industrial": ["industrial", "warehouse", "distribution center", "logistics"],
    "retail": ["retail center", "shopping center", "retail property"],
    "mixed-use": ["mixed-use", "mixed use"],
}

YEAR_BUILT_LABEL_RE = re.compile(r"year\s*built\D{0,10}((?:19|20)\d{2})", re.IGNORECASE)


def _find_first(pattern: re.Pattern, text: str) -> str | None:
    m = pattern.search(text)
    return m.group(1).strip() if m else None


def _facts_from_text(full_text: str) -> dict:
    facts: dict = {}

    address = _find_first(ADDRESS_RE, full_text)
    if address:
        facts["address"] = re.sub(r"\s+", " ", address)

    year_built = _find_first(YEAR_BUILT_LABEL_RE, full_text) or _find_first(YEAR_BUILT_RE, full_text)
    if year_built:
        facts["year_built"] = year_built

    sf = _find_first(SF_RE, full_text)
    units = _find_first(UNITS_RE, full_text)
    if sf:
        facts["sf_or_units"] = f"{sf} SF"
    elif units:
        facts["sf_or_units"] = f"{units} Units"

    occupancy = _find_first(OCCUPANCY_RE, full_text)
    if occupancy:
        facts["occupancy"] = f"{occupancy}%"

    lower_text = full_text.lower()
    for ptype, keywords in PROPERTY_TYPE_KEYWORDS.items():
        if any(kw in lower_text for kw in keywords):
            facts["property_type"] = ptype.title()
            break

    # Submarket description: grab a sentence containing "submarket".
    submarket_match = re.search(r"([^.\n]{0,220}\bsubmarket\b[^.\n]{0,220}\.)", full_text, re.IGNORECASE)
    if submarket_match:
        facts["submarket_desc"] = re.sub(r"\s+", " ", submarket_match.group(1)).strip()

    # Tenant summary: look for a sentence mentioning "tenant" with a percentage or name list.
    tenant_match = re.search(r"([^.\n]{0,220}\btenant[s]?\b[^.\n]{0,220}\.)", full_text, re.IGNORECASE)
    if tenant_match:
        facts["tenant_summary"] = re.sub(r"\s+", " ", tenant_match.group(1)).strip()

    return facts


def _text_from_pdf(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    full_text = "\n".join(page.get_text() for page in doc)
    doc.close()
    return full_text


def _text_from_pptx(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            chunks.append(cell.text)
    return "\n".join(chunks)


def extract_facts(om_path: str) -> dict:
    if om_path.lower().endswith(PPTX_EXTENSIONS):
        full_text = _text_from_pptx(om_path)
    else:
        full_text = _text_from_pdf(om_path)
    return _facts_from_text(full_text)


def _normalize_to_rgb(raw_bytes: bytes) -> Image.Image:
    """Normalize every extracted image to plain RGB regardless of source
    colorspace (CMYK, grayscale, palette, etc.), so color rendering is
    consistent across PowerPoint, browsers, and the Pillow preview
    fallback -- rather than passing through whatever colorspace the source
    document happened to embed.

    CMYK JPEGs always get an extra channel invert before the RGB convert.
    Pillow's JPEG decoder unconditionally treats every CMYK-mode JPEG as
    "Adobe inverted" and un-inverts it while decoding (see JpegImagePlugin's
    "CMYK;I" rawmode) -- but real photography embedded in OM PDFs (Acrobat
    Distiller/InDesign print pipelines, which is where a DeviceCMYK JPEG in a
    PDF almost always comes from) is not actually stored that way, so
    Pillow's blanket assumption gets it backwards and this undoes it.
    Verified against a real OM: with no extra invert, or with one gated on
    the JPEG APP14 marker's transform byte (a distinction that turned out not
    to matter -- real Adobe-pipeline exports commonly use transform 0, not
    the "documented inverted" transform 2 a previous version of this function
    checked for), photos and aerial maps both rendered as color negatives.
    """
    img = Image.open(io.BytesIO(raw_bytes))
    if img.mode == "CMYK":
        img = ImageChops.invert(img)
        img = img.convert("RGB")
    elif img.mode != "RGB":
        img = img.convert("RGB")
    return img


def _is_logo_or_chart(width: int, height: int) -> bool:
    """Heuristic filter: reject tiny icons, thin banners, and near-square
    small graphics that are almost always logos/charts/floor plans rather
    than photographs."""
    if width < 300 or height < 200:
        return True
    aspect = width / height
    if aspect > 4.0 or aspect < 0.25:
        return True  # thin banner strip
    return False


def _extract_images_from_pdf(pdf_path: str, out_dir: str) -> list[dict]:
    doc = fitz.open(pdf_path)
    candidates = []

    for page_index in range(len(doc)):
        page = doc[page_index]
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception:
                continue
            width, height = base_image.get("width", 0), base_image.get("height", 0)
            if _is_logo_or_chart(width, height):
                continue
            filename = f"page{page_index + 1}_el{xref}.png"
            out_path = str(Path(out_dir) / filename)
            try:
                normalized = _normalize_to_rgb(base_image["image"])
                normalized.save(out_path)
            except Exception:
                # A handful of embedded pictures decode to a format Pillow's
                # lazy loader can't actually render (e.g. a WMF vector image
                # PyMuPDF reports as a raster) -- the failure only surfaces
                # on save(), not on open(). Skip that one picture rather than
                # letting it take down extraction for the whole document.
                continue
            candidates.append(
                {
                    "path": out_path,
                    "width": width,
                    "height": height,
                    "page": page_index + 1,
                    "order_key": xref,
                    "area": width * height,
                }
            )

    doc.close()
    return candidates


def _extract_images_from_pptx(pptx_path: str, out_dir: str) -> list[dict]:
    prs = Presentation(pptx_path)
    candidates = []

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                width, height = image.size
            except Exception:
                continue
            if _is_logo_or_chart(width, height):
                continue
            filename = f"slide{slide_index + 1}_el{shape_index}.png"
            out_path = str(Path(out_dir) / filename)
            try:
                normalized = _normalize_to_rgb(image.blob)
                normalized.save(out_path)
            except Exception:
                # A handful of embedded pictures decode to a format Pillow's
                # lazy loader can't actually render (e.g. a WMF vector image
                # python-pptx reports as a raster) -- the failure only
                # surfaces on save(), not on open(). Skip that one picture
                # rather than letting it take down extraction for the whole
                # deck.
                continue
            candidates.append(
                {
                    "path": out_path,
                    "width": width,
                    "height": height,
                    "page": slide_index + 1,
                    "order_key": shape_index,
                    "area": width * height,
                }
            )

    return candidates


def extract_images(om_path: str, out_dir: str, max_images: int = 5) -> list[dict]:
    """Extract embedded raster images, filter out logos/charts/floor plans
    by size + aspect ratio, and rank the survivors by pixel area (bigger ==
    more likely to be a full-bleed hero photograph in an OM layout).

    Deterministic: images are always visited in the same (page, order_key)
    order a given document produces, and ties are broken by that same
    order, so the same OM always yields the same ranked list.
    """
    Path(out_dir).mkdir(parents=True, exist_ok=True)

    if om_path.lower().endswith(PPTX_EXTENSIONS):
        candidates = _extract_images_from_pptx(om_path, out_dir)
    else:
        candidates = _extract_images_from_pdf(om_path, out_dir)

    candidates.sort(key=lambda c: (-c["area"], c["page"], c["order_key"]))
    return candidates[:max_images]
