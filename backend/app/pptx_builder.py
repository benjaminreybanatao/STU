"""
Builds the DivcoWest offering summary deck.

Layout is measured from four real DW two-pagers. Two of them (400 Castro,
2390 Mission College) are PPTX files whose financial exhibit is a pasted
Excel screenshot; the other two (Landsby, Paseos) are native-vector PDFs
where the exhibit is real text, so every position, colour and row label
below is taken from direct measurement of those PDFs rather than inferred.
All four agree, so this is the firm template rather than one-off styling.

Measured specifics (slide is 960x540pt = 13.333x7.5in):
  - Theme "DW Colors 2021": green #6AA442, navy #002554, blue #0175A8,
    gray dk2 #626369. Fonts: Gandhi Sans (headings, tables), Gandhi Serif
    (body copy).
  - Page chrome on every slide: a 4.25pt green accent rule at (39.3,39.2)
    -> (89.5,39.2); the five-bar DW logo mark top-right at x 905.8-934.3
    (equal width, thickness increasing downward); an 8pt footer and a 10pt
    page number along the bottom.
  - Title: 28pt bold, "TRANSACTION SUMMARY / " in black followed by the
    property name in green.
  - Slide 1: two stacked photos at x 40.5 w 318.5 (h 176.8 each), and a
    bulleted narrative at x 402 -> 906 in 12.5pt serif, one bullet per topic.
  - Exhibit slide: two side-by-side native tables -- "Transaction Overview"
    (x 95.4-539.6) and "Sources & Uses at Close" (x 549.0-864.7) -- with
    14.6pt rows, 11.2pt text, navy header bands, #E8E8E8 section headers,
    green total rows and a blue subtotal row. Rebuilt as real pptx tables so
    the output stays editable, unlike the pasted screenshot in the sources.

Every value comes from Deal.display_value(), which walks the OM -> model ->
analyst provenance chain and falls back to PLACEHOLDER. Nothing is computed
or fabricated here -- notably the "$ Per Unit" column is only ever populated
from a per-unit figure the model states explicitly (see excel_parser), never
by dividing a total by a unit count.

Deliberately not generated: the Location Overview map, Yield Bridge chart and
Ground Lease Detail slides seen in the reference decks. Each needs
hand-assembled inputs (maps, comp sets, tenant logos, lease abstracts) that
this tool has no source for.
"""
import re

from lxml import etree
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .schema import Deal, FIELD_BY_KEY, PLACEHOLDER

# DW Colors 2021 theme, read directly out of ppt/theme/theme1.xml.
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x62, 0x63, 0x69)   # dk2
NAVY = RGBColor(0x00, 0x25, 0x54)        # accent3
GREEN = RGBColor(0x6A, 0xA4, 0x42)       # title property-name highlight + accent rule + logo
BLUE = RGBColor(0x01, 0x75, 0xA8)        # accent4 -- "Equity Subtotal" band
DARK_GREEN = RGBColor(0x3A, 0x59, 0x2A)  # defined-term highlight in body copy
SECTION_GRAY = RGBColor(0xE8, 0xE8, 0xE8)  # section / column-header rows inside the tables
ROW_ALT = RGBColor(0xF3, 0xF4, 0xF5)
BORDER_GRAY = RGBColor(0xD9, 0xDA, 0xDC)
PLACEHOLDER_RED = RGBColor(0xB0, 0x2A, 0x2A)
FOOTER_GRAY = RGBColor(0xB4, 0xB4, 0xB4)
FOOTER_GREEN = RGBColor(0x85, 0xAB, 0x49)
PAGENUM_GRAY = RGBColor(0xCA, 0xCA, 0xCA)

SANS = "Gandhi Sans"
SERIF = "Gandhi Serif"


def _pt(points: float) -> int:
    """Reference geometry is measured in points; convert to EMU."""
    return Emu(int(round(points * 12700)))


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- page chrome (measured) ---
ACCENT_RULE = (_pt(39.3), _pt(39.2), _pt(50.2), Pt(4.25))  # left, top, width, thickness
LOGO_LEFT, LOGO_WIDTH = _pt(905.8), _pt(28.5)
LOGO_BARS = [(23.9, 1.4), (28.3, 2.3), (33.6, 3.0), (39.6, 4.5), (47.1, 6.3)]  # (top pt, height pt)
FOOTER_TOP = _pt(514.0)
PAGENUM_LEFT = _pt(913.0)

TITLE_LEFT = _pt(39.3)
TITLE_TOP = _pt(55.0)
TITLE_HEIGHT = _pt(34.0)
TITLE_SIZE = Pt(28)

# --- slide 1 ---
PHOTO_LEFT, PHOTO_WIDTH = _pt(40.5), _pt(318.5)
PHOTO_TOP, PHOTO_HEIGHT, PHOTO_GAP = _pt(111.0), _pt(176.8), _pt(16.9)
BODY_LEFT, BODY_WIDTH = _pt(379.4), _pt(526.3)
BODY_TOP, BODY_HEIGHT = _pt(111.0), _pt(370.0)
BODY_SIZE = Pt(12.5)
BULLET_INDENT = Pt(22.6)
BULLET_HANG = Pt(22.6)

# --- exhibit tables ---
TABLE_TOP = _pt(120.1)
ROW_H_PT = 14.6          # reference row pitch
TABLE_SIZE_PT = 11.2     # reference body size
TABLE_BOTTOM_LIMIT_PT = 514.0 - 8.0   # keep clear of the footer

# A table row's declared height is only a *minimum*: PowerPoint grows it to fit
# the line, and it ignores wrap="none" on table cells. So a row only stays at
# the height we declare if the font's line height already fits inside it.
#
# Rather than predict PowerPoint's line height (it depends on which font gets
# substituted for Gandhi Sans on the viewer's machine, and an earlier estimate
# of 1.58x was still too low in practice), the row heights are made to sum to
# the space available and the type is then set well under that. This factor is
# deliberately pessimistic -- larger than any real font's line height plus
# residual cell padding -- so a row can't be forced to grow.
ROW_FIT_SAFETY = 1.78

# Don't let a short table balloon its rows just because there's room.
MAX_ROW_H_PT = ROW_H_PT * 1.15

# The 3-column Transaction Overview mini tables have only ~8 rows each rather
# than the single table's 24, so MAX_ROW_H_PT (tuned for that 24-row case)
# would cap them at the same tight row height -- wasting the extra vertical
# room a column now has. A row can hold 2 stacked lines (see TXN_MINI_COLS
# below), so the cap accounts for 2 lines' worth of pessimistic line-height,
# not 1, with a little headroom on top.
TXN_MINI_MAX_ROW_H_PT = TABLE_SIZE_PT * ROW_FIT_SAFETY * 2 * 1.05

# Font size is capped independently of available vertical room (which would
# otherwise push it to TABLE_SIZE_PT): a third-width column's label and value
# text was measured to need something in this range to avoid truncating on
# nearly every row -- going bigger just because there's vertical space to
# spare makes the *horizontal* fit worse, since each character gets wider.
TXN_MINI_FONT_CAP_PT = 9.5

TXN_LEFT, TXN_WIDTH = _pt(95.4), _pt(444.2)
TXN_COLS = (_pt(189.1), _pt(122.0), _pt(122.0), _pt(11.1))  # label, col2, col3, trailing spacer

# "Transaction Overview" is laid out as 3 side-by-side mini tables (label |
# value) under one shared header band, rather than one 24-row list in
# TXN_COLS's proportions. 8 rows fitting a column instead of 24 rows fitting
# the whole table means each row gets ~3x the vertical room, which is what
# actually fixes the fit -- shrinking type further within the single-table
# layout ran out of headroom.
#
# A row's Gross and $ Per Unit figures (where a row has both) go on two
# stacked lines in the value cell rather than one combined "$41,000,000 ·
# $264" line: even abbreviated, that combined string is wider than any
# third-width column at a readable size, but each figure fits fine on its
# own line -- there's ample *vertical* room in an 8-row column, just not
# horizontal. Large dollar figures (>=$1M) are abbreviated ("$210.6M") for
# the same reason; the full comma-grouped form doesn't fit at any legible
# size. A handful of the longer hand-written row/section labels are also
# abbreviated for this layout specifically (TXN_MINI_LABEL_OVERRIDES) --
# _fit_text's ellipsis truncation is a backstop for the rare case, not
# something this many rows should hit by design.
TXN_COL_GAP = _pt(6.0)
TXN_COL_WIDTH = (TXN_WIDTH - 2 * TXN_COL_GAP) // 3
TXN_MINI_LABEL_W = int(TXN_COL_WIDTH * 0.58)
TXN_MINI_COLS = (TXN_MINI_LABEL_W, TXN_COL_WIDTH - TXN_MINI_LABEL_W)

TXN_MINI_LABEL_OVERRIDES = {
    "Current Occupancy": "Occupancy",
    "Occupancy at Exit": "Occ. at Exit",
    "Projected Hold Period": "Hold Period",
    "Units at Acquisition": "Units at Acq.",
    "Total Square Feet": "Total SF",
    "Exit Price (Gross)": "Exit Price",
    "Initial Leverage (New Debt)": "Leverage",
}

SU_LEFT, SU_WIDTH = _pt(549.0), _pt(315.7)
SU_COLS = (_pt(141.0), _pt(87.0), _pt(87.7))  # label, Total, $ Per Unit

# Left "Transaction Overview" table.
# Each section: (section title | None, col2 header | None, col3 header | None, rows)
# Each row: (deck label, col2 field key | None, col3 field key | None)
TXN_OVERVIEW_SECTIONS = [
    (None, None, None, [
        ("Building Name", None, "address"),
        # Label swapped for "Total Square Feet" when the size is stated in SF --
        # see _size_row_label. The reference decks are multifamily, so their
        # own wording is the unit-count one.
        ("Units at Acquisition", None, "sf_or_units"),
        ("Current Occupancy", None, "occupancy"),
        ("Occupancy at Exit", None, "occupancy_at_exit"),
        ("Projected Hold Period", None, "hold_period"),
    ]),
    ("Pricing", "Gross", "$ Per Unit", [
        ("Purchase Price", "purchase_price", "purchase_price_per_unit"),
        ("Peak Cost", "peak_cost", "peak_cost_per_unit"),
        ("Exit Price (Gross)", "exit_price", "exit_price_per_unit"),
    ]),
    ("Cap Rate", None, None, [
        ("Year 1", None, "going_in_cap"),
        ("Market", None, "market_cap"),
        ("Exit Cap", None, "exit_cap"),
    ]),
    ("Gross Returns", None, None, [
        ("Unlevered IRR", None, "unlevered_irr"),
        ("Unlevered CFx", None, "unlevered_em"),
        ("Levered IRR", None, "levered_irr"),
        ("Levered CFx", None, "levered_em"),
    ]),
    ("Debt", None, "Gross", [
        ("Initial Leverage (New Debt)", "leverage", "gross_debt_proceeds"),
    ]),
    ("Equity", "Gross", "$ Per Unit", [
        ("Initial Equity", "initial_equity", "initial_equity_per_unit"),
        ("Peak Equity", "peak_equity", "peak_equity_per_unit"),
    ]),
]

# Right "Sources & Uses at Close" table.
# Each block: (subheader label, rows); each row: (label, total key, per-unit key, band)
SU_BLOCKS = [
    ("Total Sources", [
        ("Equity", "initial_equity", "initial_equity_per_unit", None),
        ("Gross Debt Proceeds", "gross_debt_proceeds", "gross_debt_proceeds_per_unit", None),
        ("Total Sources", "total_sources", "total_sources_per_unit", "green"),
    ]),
    ("Total Uses", [
        ("Purchase Price", "purchase_price", "purchase_price_per_unit", None),
        ("DD / Closing Costs", "dd_closing_costs", "dd_closing_costs_per_unit", None),
        ("Working Capital", "working_capital", "working_capital_per_unit", None),
        ("Equity Subtotal", "equity_subtotal", "equity_subtotal_per_unit", "blue"),
        ("Financing Cost", "financing_cost", "financing_cost_per_unit", None),
        ("Total Uses", "total_uses", "total_uses_per_unit", "green"),
    ]),
]

BAND_FILLS = {"green": GREEN, "blue": BLUE}

# Short form of PLACEHOLDER for the exhibit's narrow numeric columns, which
# don't wrap. Still unmistakably a gap, and still rendered in red.
TABLE_PLACEHOLDER = "TBD"

# The reference decks annotate the leverage figure in the Debt row as
# "65.0% LTV" rather than a bare percentage.
VALUE_SUFFIXES = {"leverage": " LTV"}

# --- asset photos grid (measured from Landsby p2) ---
GRID_COL_X = (_pt(66.4), _pt(351.5), _pt(636.7))
GRID_ROW_Y = (_pt(98.4), _pt(304.6))
GRID_CELL_W, GRID_CELL_H = _pt(270.4), _pt(178.0)


def _strip_street_suffix(street: str) -> str:
    """'400 Castro Street' -> '400 Castro' -- both source decks drop the
    street-type suffix when using the address as the property's short name."""
    stripped = re.sub(r"\b(Street|St|Avenue|Ave|Boulevard|Blvd|Road|Rd|Drive|Dr)\b\.?", "", street, flags=re.IGNORECASE)
    return re.sub(r"\s+", " ", stripped).strip()


def _short_title_label(address: str) -> str:
    """'1200 Market Street, Austin, TX 78701' -> '1200 MARKET' -- title uses
    street number + name only, dropping the suffix and city/state/zip
    (matches source decks' '400 CASTRO', '2390 MISSION COLLEGE')."""
    if address == PLACEHOLDER:
        return "PROPERTY TBD"
    return _strip_street_suffix(address.split(",")[0].strip()).upper()


def deck_filename(address: str, suffix: str = "2pager_v1") -> str:
    """Download name for the generated deck, e.g. "Landsby_2pager_v1.pptx".

    Built from the same short property name the slide titles use ("2390 Mission
    College" rather than the full postal address), so the file lands in the
    analyst's downloads named the way the firm names these decks.
    """
    name = _short_title_label(address)
    if name == "PROPERTY TBD":
        name = "Property"
    else:
        name = name.title()
    slug = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_") or "Property"
    return f"{slug}_{suffix}.pptx"


# Every table cell's left/right margin (set in _style_cell), and the total
# horizontal inset a cell's usable text width must be measured net of.
CELL_MARGIN_PT = 5.0
CELL_INSET_PT = 2 * CELL_MARGIN_PT

# Usable width of the "Building Name" value cell (col3) in the single-table
# layout -- text is measured against this so it can never wrap, which is what
# pushed a whole extra row of height into the table (see _building_name).
_BUILDING_NAME_MAX_PT = 122.0 - CELL_INSET_PT

_DEJAVU_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
_TEXT_MEASURE = ImageDraw.Draw(Image.new("RGB", (8, 8)))
_FONT_CACHE: dict[int, "ImageFont.FreeTypeFont"] = {}


def _text_width_pt(text: str, size_pt: float) -> float:
    """Rendered width of `text` at `size_pt`, measured in DejaVu Sans.

    DejaVu is metrically wider than Gandhi Sans (the deck's real font), so a
    string measured to fit here has headroom to spare in the real font too.
    """
    scale = 10  # sub-point resolution
    size = max(1, round(size_pt * scale))
    font = _FONT_CACHE.get(size)
    if font is None:
        try:
            font = ImageFont.truetype(_DEJAVU_PATH, size)
        except OSError:
            font = ImageFont.load_default()
        _FONT_CACHE[size] = font
    return _TEXT_MEASURE.textlength(text, font=font) / scale


def _fit_text(text: str, size_pt: float, max_width_pt: float) -> str:
    """Truncate `text` with an ellipsis so it's guaranteed to render on one
    line within `max_width_pt` at `size_pt` (measured per _text_width_pt).
    For any cell whose content isn't a short, hand-picked label -- dynamic
    values, or a layout with narrower columns than the label text was
    originally sized for -- so it can't silently wrap and grow its row."""
    if not text or _text_width_pt(text, size_pt) <= max_width_pt:
        return text
    t = text
    while t and _text_width_pt(t + "…", size_pt) > max_width_pt:
        t = t[:-1].rstrip()
    return (t + "…") if t else text[:1]


def _building_name(address: str, size_pt: float, max_width_pt: float = _BUILDING_NAME_MAX_PT) -> str:
    """The exhibit's "Building Name" cell takes the street portion only.

    The reference decks put a short building name there ("Landsby",
    "Paesos at Ontario"). A full postal address is roughly twice the width of
    that column, and since these cells don't wrap it would spill across the
    gap into the Sources & Uses table beside it -- or, worse, silently wrap
    onto a second line and grow that row, pushing the whole table past the
    footer. Splitting on the first comma handles "Street, City, State"
    addresses; for a bare "123 Long Street Name Boulevard" with no comma,
    _fit_text guarantees the fallback fits on one line regardless of format.
    `max_width_pt` defaults to the single-table column width but the 3-column
    layout's narrower value cell passes its own.
    """
    if address == PLACEHOLDER:
        return address
    return _fit_text(address.split(",")[0].strip(), size_pt, max_width_pt)


def _add_title(slide, address: str, prefix: str = "TRANSACTION SUMMARY / "):
    box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, Inches(10.5), TITLE_HEIGHT)
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]

    r1 = p.add_run()
    r1.text = prefix
    r1.font.size = TITLE_SIZE
    r1.font.bold = True
    r1.font.name = SANS
    r1.font.color.rgb = BLACK

    r2 = p.add_run()
    r2.text = _short_title_label(address)
    r2.font.size = TITLE_SIZE
    r2.font.bold = True
    r2.font.name = SANS
    r2.font.color.rgb = GREEN
    return box


def _add_accent_rule(slide):
    """Short green rule above the title, top-left on every reference slide."""
    left, top, width, thickness = ACCENT_RULE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(thickness)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _add_logo_bars(slide):
    """The DW mark: five equal-width green bars, thickness increasing downward."""
    for top_pt, height_pt in LOGO_BARS:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LOGO_LEFT, _pt(top_pt), LOGO_WIDTH, _pt(height_pt)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GREEN
        bar.line.fill.background()
        bar.shadow.inherit = False


def _add_footer(slide, page_no: int):
    box = slide.shapes.add_textbox(_pt(19.9), FOOTER_TOP, _pt(500), _pt(14))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    for text, color, bold in (
        ("DIVCOWEST ", FOOTER_GRAY, True),
        ("OVERVIEW MEMORANDUM ", FOOTER_GREEN, True),
        ("PRIVATE & CONFIDENTIAL", FOOTER_GRAY, False),
    ):
        run = p.add_run()
        run.text = text
        run.font.size = Pt(8)
        run.font.bold = bold
        run.font.name = SANS
        run.font.color.rgb = color

    num = slide.shapes.add_textbox(PAGENUM_LEFT, _pt(513.0), _pt(30), _pt(14))
    ntf = num.text_frame
    ntf.margin_left = 0
    ntf.margin_right = 0
    np_ = ntf.paragraphs[0]
    np_.alignment = PP_ALIGN.RIGHT
    nrun = np_.add_run()
    nrun.text = str(page_no)
    nrun.font.size = Pt(10)
    nrun.font.name = SANS
    nrun.font.color.rgb = PAGENUM_GRAY


def _new_slide(prs, address: str, page_no: int, prefix: str = "TRANSACTION SUMMARY / "):
    """A blank slide carrying the shared page chrome plus its title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_rule(slide)
    _add_logo_bars(slide)
    _add_title(slide, address, prefix=prefix)
    _add_footer(slide, page_no)
    return slide


def _split_address(address: str) -> tuple[str, str]:
    """'1200 Market Street, Austin, TX 78701' -> ('1200 Market Street', 'Austin, TX')."""
    if address == PLACEHOLDER:
        return address, address
    parts = [p.strip() for p in address.split(",")]
    street = _strip_street_suffix(parts[0]) if parts else address
    city_state = ", ".join(parts[1:-1]) if len(parts) > 2 else (parts[1] if len(parts) > 1 else "")
    return street, city_state


def _build_narrative_paragraphs(deal: Deal) -> list[str]:
    """One short paragraph per topic (property, ownership/description,
    tenancy, pricing, business plan, capital structure, exit & returns),
    matching the source decks' pattern of separate sentence-groups rather
    than a single merged block."""
    dv = deal.display_value
    paras = []

    ptype = dv("property_type")
    ptype_phrase = ptype.lower() if ptype != PLACEHOLDER else ptype
    article = "an" if ptype_phrase[:1].lower() in "aeiou" else "a"
    street, city_state = _split_address(dv("address"))
    location = f" located in {city_state}" if city_state else ""
    paras.append(
        f'Opportunity to acquire {street}, {article} {ptype_phrase} building totaling '
        f'{dv("sf_or_units")}{location} (the “Investment” or the “Property”).'
    )

    year_built = dv("year_built")
    if year_built != PLACEHOLDER:
        paras.append(f"The Property was built in {year_built}.")
    submarket = dv("submarket_desc")
    if submarket != PLACEHOLDER:
        paras.append(submarket)

    tenant_summary = dv("tenant_summary")
    occupancy = dv("occupancy")
    walt = dv("walt")
    tenancy_sentence = f"The Property is {occupancy} occupied"
    if tenant_summary != PLACEHOLDER:
        tenancy_sentence += f"; {tenant_summary}"
    if walt != PLACEHOLDER:
        tenancy_sentence += f", resulting in a {walt} WALT"
    paras.append(tenancy_sentence + ".")

    paras.append(
        f'The projected purchase price for the building is {dv("purchase_price")} '
        f'({dv("purchase_price_per_unit")}), resulting in a {dv("going_in_cap")} going-in cap rate.'
    )

    lease_term = dv("lease_term_assumption")
    downtime = dv("downtime_assumption")
    exit_assumption = dv("exit_assumption")
    if any(v != PLACEHOLDER for v in (lease_term, downtime, exit_assumption)):
        bits = []
        if lease_term != PLACEHOLDER:
            bits.append(f"a {lease_term} lease term")
        if downtime != PLACEHOLDER:
            bits.append(f"{downtime} downtime")
        if exit_assumption != PLACEHOLDER:
            bits.append(exit_assumption.rstrip("."))
        paras.append("Base case business plan assumes " + "; ".join(bits) + ".")

    peak_equity = dv("peak_equity")
    debt_sentence = f'With {dv("leverage")} leverage at {dv("debt_rate")}, the initial equity is {dv("initial_equity")}'
    if peak_equity != PLACEHOLDER:
        debt_sentence += f" and the projected peak equity for the transaction is {peak_equity}"
    paras.append(debt_sentence + ".")

    paras.append(
        f'Assuming the Property is sold in {dv("hold_period")} at a {dv("exit_cap")} exit cap rate, '
        f'projected unlevered returns are {dv("unlevered_irr")} / {dv("unlevered_em")} and levered '
        f'returns are {dv("levered_irr")} / {dv("levered_em")}.'
    )

    return paras


def _set_bullet(paragraph, char="•", color=GREEN):
    """python-pptx has no high-level bullet API -- set marL/indent plus a
    buFont/buChar pair directly on the paragraph's pPr, in the schema order
    CT_TextParagraphProperties requires (after spcAft, before defRPr)."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(BULLET_INDENT)))
    pPr.set("indent", str(-int(BULLET_HANG)))

    buClr = etree.SubElement(pPr, qn("a:buClr"))
    solidFill = etree.SubElement(buClr, qn("a:srgbClr"))
    solidFill.set("val", str(color))

    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")

    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)


def _render_narrative(slide, deal: Deal, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0

    paragraphs = _build_narrative_paragraphs(deal)
    first = True
    for para_text in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        _set_bullet(p)

        # Highlight the defined terms in the opening sentence, and any
        # placeholder text everywhere else, distinctly -- everything else is
        # plain body copy.
        segments = re.split(r'(“Investment”|“Property”|' + re.escape(PLACEHOLDER) + r')', para_text)
        for seg in segments:
            if not seg:
                continue
            run = p.add_run()
            run.text = seg
            run.font.size = BODY_SIZE
            run.font.name = SERIF
            if seg == PLACEHOLDER:
                run.font.bold = True
                run.font.color.rgb = PLACEHOLDER_RED
            elif seg in ('“Investment”', '“Property”'):
                run.font.color.rgb = DARK_GREEN
            else:
                run.font.color.rgb = BLACK
    return box


def _add_photos(slide, image_paths, left, top, width, height, gap=Inches(0.08)):
    slots = 2
    slot_h = Emu(int((int(height) - int(gap)) / slots))
    paths = (image_paths or [None, None])[:slots]
    while len(paths) < slots:
        paths.append(paths[-1] if paths else None)

    for i, path in enumerate(paths):
        slot_top = Emu(int(top) + i * (int(slot_h) + int(gap)))
        if not path:
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, slot_top, width, slot_h)
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = ROW_ALT
            placeholder.line.color.rgb = BORDER_GRAY
            tf = placeholder.text_frame
            tf.text = "No property photo extracted from OM"
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.italic = True
            p.runs[0].font.size = Pt(10)
            p.runs[0].font.color.rgb = GRAY_DARK
            continue

        pic = slide.shapes.add_picture(path, left, slot_top, width=width, height=slot_h)
        native_w, native_h = pic.image.size
        target_ratio = width / slot_h
        native_ratio = native_w / native_h
        if native_ratio > target_ratio:
            crop = (1 - target_ratio / native_ratio) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (1 - native_ratio / target_ratio) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop


def _build_slide1(prs, deal: Deal, hero_image_paths, page_no: int):
    slide = _new_slide(prs, deal.display_value("address"), page_no)
    total_photo_h = Emu(int(PHOTO_HEIGHT) * 2 + int(PHOTO_GAP))
    _add_photos(slide, hero_image_paths, PHOTO_LEFT, PHOTO_TOP, PHOTO_WIDTH, total_photo_h, gap=PHOTO_GAP)
    _render_narrative(slide, deal, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT)
    return slide


def _size_row_label(deal: Deal, default: str) -> str:
    """Keep the size row's label honest: an SF figure must not sit under a
    row captioned "Units at Acquisition"."""
    value = deal.display_value("sf_or_units")
    if value != PLACEHOLDER and re.search(r"\b(SF|RSF|square feet)\b", value, re.IGNORECASE):
        return "Total Square Feet"
    return default


def _cell_value(deal: Deal, key: str | None) -> tuple[str, bool]:
    """(text, is_placeholder) for a table cell, with any unit annotation the
    reference decks add.

    Missing values collapse to the short TABLE_PLACEHOLDER: the full
    "TBD -- confirm with sponsor" wording is ~1.4x the width of these numeric
    columns and they don't wrap, so the long form would overlap its
    neighbours. The full wording still appears in the slide 1 narrative and in
    the gap-analysis panel; here the red short form carries the same meaning.
    """
    if not key:
        return "", False
    value = deal.display_value(key)
    if value == PLACEHOLDER:
        return TABLE_PLACEHOLDER, True
    return value + VALUE_SUFFIXES.get(key, ""), False


_CURRENCY_RE = re.compile(r"^\$([\d,]+)$")


def _compact_currency(text: str) -> str:
    """"$210,559,365" -> "$210.6M" -- used only in the 3-column layout's
    narrow value cells, where a figure at this magnitude doesn't fit the
    full comma-grouped form at any legible size. Left alone below $1M (a
    $ Per Unit figure, say) and for anything that isn't a bare dollar amount
    (percentages, multiples, "TBD", SF, ...)."""
    m = _CURRENCY_RE.match(text)
    if not m:
        return text
    n = int(m.group(1).replace(",", ""))
    if abs(n) >= 1_000_000_000:
        return f"${n / 1_000_000_000:.1f}B"
    if abs(n) >= 1_000_000:
        return f"${n / 1_000_000:.1f}M"
    return text


_BORDER_TAGS = ("a:lnL", "a:lnR", "a:lnT", "a:lnB")


def _clear_cell_borders(cell):
    """The reference exhibit has no gridlines -- only fills separate the rows.
    python-pptx's default table style draws borders, so switch them off."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in _BORDER_TAGS:
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(tag))
        etree.SubElement(ln, qn("a:noFill"))


def _style_cell(cell, text, *, bold=False, color=BLACK, align=PP_ALIGN.LEFT, fill=None,
                is_placeholder=False, size_pt=TABLE_SIZE_PT):
    """Write one table cell (one line). Placeholders render red so missing
    data is obvious on the deck itself."""
    _style_cell_lines(cell, [(text, is_placeholder)], bold=bold, color=color, align=align,
                       fill=fill, size_pt=size_pt)


def _style_cell_lines(cell, lines: list[tuple[str, bool]], *, bold=False, color=BLACK,
                       align=PP_ALIGN.LEFT, fill=None, size_pt=TABLE_SIZE_PT):
    """Write one or more stacked lines into a table cell, each its own
    paragraph. The 3-column Transaction Overview layout uses this for a
    row's Gross and $ Per Unit figures: no combined single line fits a
    third-width column, but two short lines stacked vertically do, and a
    mini table's rows have vertical room to spare (see TXN_MINI_MAX_ROW_H_PT)."""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    _clear_cell_borders(cell)
    cell.margin_left = Pt(CELL_MARGIN_PT)
    cell.margin_right = Pt(CELL_MARGIN_PT)
    cell.margin_top = 0
    cell.margin_bottom = 0
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = cell.text_frame
    tf.word_wrap = False
    on_dark_band = fill in (NAVY, GREEN, BLUE)
    for i, (text, is_placeholder) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # Leaving these unset means PowerPoint fills them in from the
        # theme's default text style at render time -- which can add
        # several points of paragraph spacing per row that our row-height
        # safety margin (based purely on font line-height) never accounted
        # for, on top of the font substitution it does account for. 24 rows
        # of a few unaccounted points each is enough to visibly push the
        # table past the footer even though every row individually looks
        # fine. Pin all three explicitly so a row's rendered height is
        # driven only by the font size we set.
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = SANS
        # Placeholders normally go red to flag the gap, but red on a dark
        # band is unreadable -- there the passed-in colour (white) already
        # stands out.
        run.font.color.rgb = PLACEHOLDER_RED if (is_placeholder and not on_dark_band) else color


def _split_balanced(items: list, n: int) -> list[list]:
    """Split `items` into `n` contiguous, size-balanced chunks (earlier
    chunks get the extra items if it doesn't divide evenly), for tiling a
    flat list into side-by-side columns."""
    base, rem = divmod(len(items), n)
    sizes = [base + 1 if i < rem else base for i in range(n)]
    chunks = []
    start = 0
    for size in sizes:
        chunks.append(items[start:start + size])
        start += size
    return chunks


def _txn_mini_value(deal: Deal, c2: str | None, c3: str | None, size_pt: float,
                     max_width_pt: float) -> list[tuple[str, bool]]:
    """A row's value cell lines for the 3-column layout: Gross/primary (c2)
    and $ Per Unit/secondary (c3) figures each get their own line rather
    than combining into one -- no combined single line fits a third-width
    column at any legible size, but each figure fits fine on its own, and a
    mini table's rows have vertical room to spare. Falls back to one line
    when the row only has one of the two."""
    v2, ph2 = _cell_value(deal, c2)
    v3, ph3 = _cell_value(deal, c3)
    if c3 == "address" and not ph3:
        return [(_building_name(v3, size_pt, max_width_pt=max_width_pt), ph3)]
    lines = [(_compact_currency(v), ph) for v, ph in ((v2, ph2), (v3, ph3)) if v]
    if not lines:
        return [(v2, ph2)]
    return [(_fit_text(text, size_pt, max_width_pt), ph) for text, ph in lines]


def _build_txn_overview_table_3col(slide, deal: Deal):
    """"Transaction Overview" as 3 side-by-side mini tables under one shared
    header band -- see the comment on TXN_COL_GAP for why."""
    flat_rows = []
    for title, c2h, c3h, section_rows in TXN_OVERVIEW_SECTIONS:
        if title is not None:
            flat_rows.append(("section", title))
        for label, c2, c3 in section_rows:
            flat_rows.append(("data", label, c2, c3))

    columns = _split_balanced(flat_rows, 3)
    max_rows = max(len(c) for c in columns)

    header_size_pt, header_row_h = _table_metrics(1)
    _, header_table = _new_table(slide, TXN_LEFT, TABLE_TOP, TXN_WIDTH, (TXN_WIDTH,), 1, header_row_h)
    _fill_header_band(header_table, 0, "Transaction Overview", 1, header_size_pt)

    body_top_pt = TABLE_TOP / 12700 + header_row_h / 12700
    size_pt, row_h = _table_metrics(max_rows, top_pt=body_top_pt, max_row_pt=TXN_MINI_MAX_ROW_H_PT)
    size_pt = min(size_pt, TXN_MINI_FONT_CAP_PT)
    label_w, value_w = TXN_MINI_COLS
    label_fit_pt = label_w / 12700 - CELL_INSET_PT
    value_fit_pt = value_w / 12700 - CELL_INSET_PT
    section_fit_pt = TXN_COL_WIDTH / 12700 - CELL_INSET_PT

    for i, column_rows in enumerate(columns):
        left = TXN_LEFT + i * (TXN_COL_WIDTH + TXN_COL_GAP)
        _, table = _new_table(slide, left, _pt(body_top_pt), TXN_COL_WIDTH,
                               (label_w, value_w), len(column_rows), row_h)
        for r, row in enumerate(column_rows):
            if row[0] == "section":
                _, title = row
                for col in range(2):
                    _style_cell(table.cell(r, col), "", fill=SECTION_GRAY, size_pt=size_pt)
                origin = table.cell(r, 0)
                origin.merge(table.cell(r, 1))
                _style_cell(origin, _fit_text(title, size_pt, section_fit_pt),
                            bold=True, fill=SECTION_GRAY, size_pt=size_pt)
            else:
                _, label, c2, c3 = row
                if c3 == "sf_or_units":
                    label = _size_row_label(deal, label)
                label = TXN_MINI_LABEL_OVERRIDES.get(label, label)
                lines = _txn_mini_value(deal, c2, c3, size_pt, value_fit_pt)
                _style_cell(table.cell(r, 0), _fit_text(label, size_pt, label_fit_pt), size_pt=size_pt)
                _style_cell_lines(table.cell(r, 1), lines, align=PP_ALIGN.CENTER, size_pt=size_pt)


def _su_row_count() -> int:
    # header + per block: optional spacer + subheader + its rows
    return 1 + sum(
        (1 if i else 0) + 1 + len(block_rows)
        for i, (_sub, block_rows) in enumerate(SU_BLOCKS)
    )


def _table_metrics(n_rows: int, top_pt: float | None = None,
                    max_row_pt: float = MAX_ROW_H_PT) -> tuple[float, int]:
    """(font size in pt, row height in EMU) for a table of `n_rows` rows
    starting at `top_pt` (defaults to the shared exhibit table top).

    The rows are sized to fill the space between the table top and the footer,
    so the table's total height is fixed by construction rather than by what
    PowerPoint decides each row needs. The type is then set far enough under
    that height that no font substitution can push a row taller and start
    cascading the table off the slide.
    """
    if top_pt is None:
        top_pt = TABLE_TOP / 12700
    available_pt = TABLE_BOTTOM_LIMIT_PT - top_pt
    row_pt = min(max_row_pt, available_pt / n_rows)
    font_pt = min(TABLE_SIZE_PT, row_pt / ROW_FIT_SAFETY)
    return font_pt, _pt(row_pt)


def _new_table(slide, left, top, width, col_widths, n_rows, row_h):
    shape = slide.shapes.add_table(
        n_rows, len(col_widths), left, top, width, Emu(int(row_h) * n_rows)
    )
    table = shape.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw
    for r in range(n_rows):
        table.rows[r].height = row_h
    # python-pptx tables default to a banded first-row/column style; turn the
    # theming off so our explicit fills are what actually shows.
    tbl = table._tbl.tblPr
    tbl.set("firstRow", "0")
    tbl.set("bandRow", "0")
    return shape, table


def _fill_header_band(table, row: int, title: str, n_cols: int, size_pt: float):
    """Navy band spanning the table with the block title centred across it."""
    for c in range(n_cols):
        _style_cell(table.cell(row, c), "", fill=NAVY, size_pt=size_pt)
    origin = table.cell(row, 0)
    origin.merge(table.cell(row, n_cols - 1))
    _style_cell(origin, title, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY,
                size_pt=size_pt)


def _build_sources_uses_table(slide, deal: Deal, size_pt: float, row_h: int):
    n_cols = len(SU_COLS)
    rows = [("header", "Sources & Uses at Close", None, None, None)]
    for i, (subheader, block_rows) in enumerate(SU_BLOCKS):
        if i:
            rows.append(("spacer", "", None, None, None))
        rows.append(("section", subheader, "Total", "$ Per Unit", None))
        for label, total_key, pu_key, band in block_rows:
            rows.append(("data", label, total_key, pu_key, band))

    _, table = _new_table(slide, SU_LEFT, TABLE_TOP, SU_WIDTH, SU_COLS, len(rows), row_h)

    for r, (kind, a, b, c, band) in enumerate(rows):
        if kind == "header":
            _fill_header_band(table, r, a, n_cols, size_pt)
        elif kind == "spacer":
            for col in range(n_cols):
                _style_cell(table.cell(r, col), "", size_pt=size_pt)
        elif kind == "section":
            for col in range(n_cols):
                _style_cell(table.cell(r, col), "", fill=SECTION_GRAY, size_pt=size_pt)
            _style_cell(table.cell(r, 0), a, bold=True, fill=SECTION_GRAY, size_pt=size_pt)
            _style_cell(table.cell(r, 1), b, bold=True, align=PP_ALIGN.CENTER, fill=SECTION_GRAY, size_pt=size_pt)
            _style_cell(table.cell(r, 2), c, bold=True, align=PP_ALIGN.CENTER, fill=SECTION_GRAY, size_pt=size_pt)
        else:
            fill = BAND_FILLS.get(band)
            color = WHITE if fill is not None else BLACK
            bold = fill is not None
            v1, ph1 = _cell_value(deal, b)
            v2, ph2 = _cell_value(deal, c)
            _style_cell(table.cell(r, 0), a, bold=bold, color=color, fill=fill, size_pt=size_pt)
            _style_cell(table.cell(r, 1), v1, bold=bold, color=color,
                        align=PP_ALIGN.RIGHT, fill=fill, is_placeholder=ph1, size_pt=size_pt)
            _style_cell(table.cell(r, 2), v2, bold=bold, color=color,
                        align=PP_ALIGN.RIGHT, fill=fill, is_placeholder=ph2, size_pt=size_pt)


def _build_exhibit_slide(prs, deal: Deal, address: str, page_no: int, case_label: str | None = None):
    prefix = "TRANSACTION OVERVIEW / " if case_label else "TRANSACTION SUMMARY / "
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_rule(slide)
    _add_logo_bars(slide)
    _add_footer(slide, page_no)

    # With a named case, the reference decks put the case in the green slot
    # ("TRANSACTION OVERVIEW / 10Y Hold"); otherwise the property name goes there.
    box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, Inches(10.5), TITLE_HEIGHT)
    tf = box.text_frame
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = prefix
    r1.font.size = TITLE_SIZE
    r1.font.bold = True
    r1.font.name = SANS
    r1.font.color.rgb = BLACK
    r2 = p.add_run()
    r2.text = case_label if case_label else _short_title_label(address)
    r2.font.size = TITLE_SIZE
    r2.font.bold = True
    r2.font.name = SANS
    r2.font.color.rgb = GREEN

    # Transaction Overview is sized independently now that it's 3 mini
    # tables of ~8 rows each rather than one 24-row table -- its type would
    # otherwise be forced down to Sources & Uses's much smaller size for no
    # reason.
    _build_txn_overview_table_3col(slide, deal)
    su_size_pt, su_row_h = _table_metrics(_su_row_count())
    _build_sources_uses_table(slide, deal, su_size_pt, su_row_h)
    return slide


def _build_asset_photos_slide(prs, deal: Deal, image_paths, page_no: int):
    """3x2 grid of OM photos, matching the reference "ASSET PHOTOS" slide."""
    slide = _new_slide(prs, deal.display_value("address"), page_no, prefix="ASSET PHOTOS / ")
    for i, path in enumerate(image_paths[:6]):
        left = GRID_COL_X[i % 3]
        top = GRID_ROW_Y[i // 3]
        pic = slide.shapes.add_picture(path, left, top, width=GRID_CELL_W, height=GRID_CELL_H)
        native_w, native_h = pic.image.size
        target_ratio = GRID_CELL_W / GRID_CELL_H
        native_ratio = native_w / native_h
        if native_ratio > target_ratio:
            crop = (1 - target_ratio / native_ratio) / 2
            pic.crop_left = pic.crop_right = crop
        else:
            crop = (1 - native_ratio / target_ratio) / 2
            pic.crop_top = pic.crop_bottom = crop
    return slide


def build_deck(deal: Deal, hero_image_path, output_path: str, downside_deal: Deal | None = None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    image_paths = hero_image_path if isinstance(hero_image_path, list) else [hero_image_path]
    image_paths = [p for p in image_paths if p]
    address = deal.display_value("address")

    page = 0
    _build_slide1(prs, deal, image_paths, page)

    page += 1
    _build_exhibit_slide(prs, deal, address, page,
                         case_label="Base Case" if downside_deal is not None else None)

    if downside_deal is not None:
        page += 1
        _build_exhibit_slide(prs, downside_deal, address, page, case_label="Downside Case")

    if len(image_paths) >= 3:
        page += 1
        _build_asset_photos_slide(prs, deal, image_paths, page)

    prs.save(output_path)
    return output_path
