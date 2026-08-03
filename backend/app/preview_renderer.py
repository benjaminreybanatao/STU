"""
Native pptx -> PNG preview renderer.

The primary preview path is LibreOffice headless (soffice --convert-to pdf,
then rasterize with PyMuPDF) -- see main.py. That gives a pixel-accurate
render when a working LibreOffice install is available. This module is the
fallback used when soffice is missing, broken, or times out in a given
deployment environment: it walks the saved .pptx directly with python-pptx
and redraws each slide's shapes (rectangles, text, pictures, tables) onto a
Pillow canvas. It's an approximation, not a pixel-perfect render, but it
means the "preview in browser" feature degrades gracefully instead of
failing outright when the host has no (working) office suite installed.
"""
import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

TARGET_WIDTH_PX = 1600
FONT_PATH_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_PATH_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

_font_cache: dict = {}


def _font(size_pt: float, bold: bool) -> ImageFont.FreeTypeFont:
    key = (round(size_pt), bold)
    if key not in _font_cache:
        path = FONT_PATH_BOLD if bold and Path(FONT_PATH_BOLD).exists() else FONT_PATH_REGULAR
        try:
            _font_cache[key] = ImageFont.truetype(path, max(int(size_pt), 6))
        except Exception:
            _font_cache[key] = ImageFont.load_default()
    return _font_cache[key]


def _rgb(color_format):
    try:
        if color_format and color_format.type is not None:
            rgb = color_format.rgb
            return (rgb[0], rgb[1], rgb[2]) if rgb else (0, 0, 0)
    except Exception:
        pass
    return (0, 0, 0)


def _wrap_text(draw, text, font, max_width_px):
    words = text.split()
    lines, current = [], ""
    for word in words:
        trial = f"{current} {word}".strip()
        if draw.textlength(trial, font=font) <= max_width_px or not current:
            current = trial
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def _bullet_info(para):
    """Read marL/indent/buChar/buClr directly off pPr -- python-pptx has no
    high-level bullet API. Returns (char, marL_emu, hang_emu, color) or None."""
    pPr = para._p.find(qn("a:pPr"))
    if pPr is None:
        return None
    buChar = pPr.find(qn("a:buChar"))
    if buChar is None:
        return None
    marL = int(pPr.get("marL") or 0)
    indent = int(pPr.get("indent") or 0)
    buClr = pPr.find(qn("a:buClr"))
    color = (50, 55, 65)
    if buClr is not None:
        srgb = buClr.find(qn("a:srgbClr"))
        if srgb is not None:
            v = srgb.get("val")
            color = (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    return buChar.get("char", "•"), marL, indent, color


def _draw_text_frame(draw, shape, scale, box):
    left, top, width, height = box
    tf = shape.text_frame
    y = top + 4
    for para in tf.paragraphs:
        runs = para.runs or []
        text = "".join(r.text for r in runs)
        if not text.strip():
            y += 14
            continue
        size_pt = 12
        bold = False
        color = (50, 55, 65)
        if runs:
            r0 = runs[0]
            if r0.font.size:
                size_pt = r0.font.size.pt
            bold = bool(r0.font.bold)
            color = _rgb(r0.font.color)
        font = _font(size_pt * scale * 1.15, bold)
        align = getattr(para, "alignment", None)

        bullet = _bullet_info(para)
        text_indent_px = 0
        if bullet:
            char, marL_emu, indent_emu, bullet_color = bullet
            text_indent_px = marL_emu * scale
            bullet_x = left + text_indent_px + indent_emu * scale
            bullet_font = _font(size_pt * scale * 1.15, False)
            draw.text((bullet_x, y), char, font=bullet_font, fill=bullet_color)

        lines = _wrap_text(draw, text, font, width - 8 - text_indent_px)

        # Single-run paragraphs wrap normally. Multi-run paragraphs that fit on
        # one line are drawn run-by-run so per-run colour survives (the two-tone
        # title, green defined terms, red placeholders); longer mixed-run
        # paragraphs fall back to wrapping in the first run's style.
        if len(runs) > 1 and len(lines) == 1:
            _draw_runs(draw, runs, left + text_indent_px, y, width - text_indent_px,
                       font.size, scale, align, default_size=size_pt)
            y += font.size + 4
            continue

        for line in lines:
            line_w = draw.textlength(line, font=font)
            if align is not None and str(align).endswith("CENTER (2)"):
                x = left + max((width - line_w) / 2, 0)
            elif align is not None and "RIGHT" in str(align):
                x = left + max(width - line_w - 4, 0)
            else:
                x = left + 4 + text_indent_px
            draw.text((x, y), line, font=font, fill=color)
            y += font.size + 4


def _draw_picture(canvas, shape, scale, box):
    left, top, width, height = box
    try:
        blob = shape.image.blob
        img = Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception:
        return
    cl, ct, cr, cb = (shape.crop_left or 0, shape.crop_top or 0, shape.crop_right or 0, shape.crop_bottom or 0)
    w, h = img.size
    box_img = (int(w * cl), int(h * ct), int(w * (1 - cr)), int(h * (1 - cb)))
    if box_img[2] > box_img[0] and box_img[3] > box_img[1]:
        img = img.crop(box_img)
    img = img.resize((max(int(width), 1), max(int(height), 1)))
    canvas.paste(img, (int(left), int(top)))


def _cell_fill(cell):
    """Solid fill colour of a table cell, or None when it has no solid fill.

    Guarding on `fill.type is not None` isn't enough: an explicitly
    background-filled cell has a type but no foreground colour, and asking for
    one raises.
    """
    try:
        if cell.fill.type == MSO_FILL.SOLID:
            return _rgb(cell.fill.fore_color)
    except Exception:
        pass
    return None


def _draw_table(canvas, draw, shape, scale, box):
    table = shape.table
    left, top, _, _ = box
    col_widths = [Emu(c.width).emu * scale for c in table.columns]
    row_heights = [Emu(r.height).emu * scale for r in table.rows]

    y = top
    for r_idx, row in enumerate(table.rows):
        x = left
        for c_idx, cell in enumerate(row.cells):
            cw, rh = col_widths[c_idx], row_heights[r_idx]
            # Cells covered by a merge repeat the origin cell's content; drawing
            # them would double up the text and re-stroke the interior borders.
            if getattr(cell, "is_spanned", False):
                x += cw
                continue
            if getattr(cell, "span_width", 1) > 1:
                cw = sum(col_widths[c_idx:c_idx + cell.span_width])
            if getattr(cell, "span_height", 1) > 1:
                rh = sum(row_heights[r_idx:r_idx + cell.span_height])

            # No outline: the decks we render switch cell borders off and
            # separate rows with fills alone, so stroking here would invent
            # gridlines that aren't in the real file.
            fill = _cell_fill(cell)
            if fill is not None:
                draw.rectangle([x, y, x + cw, y + rh], fill=fill)

            para = cell.text_frame.paragraphs[0]
            runs = [r for p in cell.text_frame.paragraphs for r in p.runs if r.text]
            if runs:
                _draw_runs(draw, runs, x, y, cw, rh, scale, para.alignment, default_size=11)
            x += col_widths[c_idx]
        y += row_heights[r_idx]


def _draw_runs(draw, runs, x, y, width, height, scale, align, default_size=11):
    """Draw a cell's runs on one baseline, honouring each run's own styling.

    Styling per run matters here: the exhibit's placeholder cells colour just
    the `TBD` text red, and the slide titles are two-tone.
    """
    pieces = []
    for run in runs:
        size_pt = run.font.size.pt if run.font.size else default_size
        font = _font(size_pt * scale * 1.15, bool(run.font.bold))
        pieces.append((run.text, font, _rgb(run.font.color)))

    total_w = sum(draw.textlength(t, font=f) for t, f, _ in pieces)
    tallest = max(f.size for _, f, _ in pieces)

    if align is not None and "RIGHT" in str(align):
        tx = x + width - total_w - 6
    elif align is not None and str(align).endswith("CENTER (2)"):
        tx = x + max((width - total_w) / 2, 0)
    else:
        tx = x + 6

    ty = y + height / 2 - tallest / 2
    for text, font, color in pieces:
        draw.text((tx, ty), text, font=font, fill=color)
        tx += draw.textlength(text, font=font)


def _draw_shape(canvas, draw, shape, scale):
    box = (shape.left * scale, shape.top * scale, shape.width * scale, shape.height * scale)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        _draw_picture(canvas, shape, scale, box)
        return

    if shape.has_table:
        _draw_table(canvas, draw, shape, scale, box)
        return

    try:
        if shape.fill.type is not None and shape.fill.type == 1:  # solid
            fill_color = _rgb(shape.fill.fore_color)
            left, top, width, height = box
            draw.rectangle([left, top, left + width, top + height], fill=fill_color)
    except Exception:
        pass

    if shape.has_text_frame and shape.text_frame.text.strip():
        _draw_text_frame(draw, shape, scale, box)


def render_pptx_to_png(pptx_path: str, out_dir: str) -> list[str]:
    prs = Presentation(pptx_path)
    scale = TARGET_WIDTH_PX / prs.slide_width
    target_h = int(prs.slide_height * scale)

    Path(out_dir).mkdir(parents=True, exist_ok=True)
    out_paths = []
    for i, slide in enumerate(prs.slides, start=1):
        canvas = Image.new("RGB", (TARGET_WIDTH_PX, target_h), (255, 255, 255))
        draw = ImageDraw.Draw(canvas)
        for shape in slide.shapes:
            _draw_shape(canvas, draw, shape, scale)
        out_path = str(Path(out_dir) / f"slide-{i:02d}.png")
        canvas.save(out_path)
        out_paths.append(out_path)
    return out_paths
